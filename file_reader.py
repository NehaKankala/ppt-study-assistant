# file_reader.py
from pptx import Presentation
import fitz  # PyMuPDF

def extract_slide_texts(ppt_file):
    prs = Presentation(ppt_file)
    slide_texts = []
    for slide in prs.slides:
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
        slide_texts.append(text.strip())
    return slide_texts

def extract_pdf_texts(pdf_file):
    doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
    texts = []
    for page in doc:
        texts.append(page.get_text())
    return texts

