from pptx import Presentation

def extract_slide_texts(ppt_file):
    prs = Presentation(ppt_file)
    slide_texts = []
    for slide in prs.slides:
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
        slide_texts.append(text.strip())
    return slide_texts

